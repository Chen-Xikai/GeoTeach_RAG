"""
GeoTeach AI Agent - 文档处理模块

提供文件读取和文档加载功能，支持多模态处理。
"""

import os
from pathlib import Path
from typing import List, Dict, Optional


# 支持的文件扩展名
SUPPORTED_EXTENSIONS = {'.pdf', '.docx', '.txt', '.md', '.pptx'}


def read_file(filepath: str) -> str:
    """读取文件内容"""
    ext = Path(filepath).suffix.lower()
    
    if ext == '.pptx':
        return _read_pptx(filepath)
    elif ext == '.pdf':
        return _read_pdf(filepath)
    elif ext == '.docx':
        return _read_docx(filepath)
    elif ext in ['.txt', '.md']:
        return _read_text(filepath)
    else:
        raise ValueError(f"不支持的文件格式: {ext}")


def _read_pptx(filepath: str) -> str:
    """读取PPT文件"""
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("python-pptx未安装，请运行: pip install python-pptx")
    
    prs = Presentation(filepath)
    texts = []
    
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = [f"--- 第{slide_num}页 ---"]
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text)
            
            # 提取表格
            if shape.has_table:
                table = shape.table
                rows = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    rows.append(" | ".join(cells))
                slide_texts.append("\n".join(rows))
        
        if len(slide_texts) > 1:  # 有内容
            texts.append("\n".join(slide_texts))
    
    return "\n\n".join(texts)


def _read_pdf(filepath: str) -> str:
    """读取PDF文件"""
    try:
        from pypdf import PdfReader
        reader = PdfReader(filepath)
        texts = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                texts.append(text)
        return "\n\n".join(texts)
    except Exception as e:
        raise ValueError(f"PDF读取失败: {e}")


def _read_docx(filepath: str) -> str:
    """读取DOCX文件"""
    try:
        from docx import Document
        doc = Document(filepath)
        
        texts = []
        for para in doc.paragraphs:
            if para.text.strip():
                texts.append(para.text)
        
        # 提取表格
        for table in doc.tables:
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(" | ".join(cells))
            texts.append("\n".join(rows))
        
        return "\n".join(texts)
    except Exception as e:
        raise ValueError(f"DOCX读取失败: {e}")


def _read_text(filepath: str) -> str:
    """读取文本文件"""
    try:
        with open(filepath, 'r', encoding='utf-8') as f:
            return f.read()
    except UnicodeDecodeError:
        with open(filepath, 'r', encoding='gbk') as f:
            return f.read()


def load_single_document(filepath: str, use_multimodal: bool = False) -> List[dict]:
    """
    加载单个文档
    
    Args:
        filepath: 文件路径
        use_multimodal: 是否使用多模态处理（提取图片、OCR等）
    
    Returns:
        文档列表
    """
    filepath = str(filepath)
    ext = Path(filepath).suffix.lower()
    
    if ext not in SUPPORTED_EXTENSIONS:
        raise ValueError(f"不支持的文件格式: {ext}")
    
    # 读取文本内容
    content = read_file(filepath)
    
    if not content.strip():
        return []
    
    # 构建文档对象
    doc = {
        "page_content": content,
        "metadata": {
            "source": filepath,
            "filename": Path(filepath).name,
            "file_type": ext,
        }
    }
    
    # 多模态处理（提取图片）
    if use_multimodal and ext in ['.pptx', '.pdf', '.docx']:
        try:
            from core.multimodal import MultimodalProcessor
            processor = MultimodalProcessor(use_vision_api=True)
            docs = processor.process_document(filepath)
            if docs:
                return docs
        except Exception as e:
            print(f"多模态处理失败，使用普通模式: {e}")
    
    return [doc]


def load_all_documents(*dirs: Path, use_multimodal: bool = False) -> List[dict]:
    """加载所有目录下的文档"""
    documents = []
    
    for dir_path in dirs:
        if not dir_path.exists():
            continue
        
        for ext in SUPPORTED_EXTENSIONS:
            files = list(dir_path.rglob(f"*{ext}"))
            for file_path in files:
                try:
                    docs = load_single_document(str(file_path), use_multimodal=use_multimodal)
                    documents.extend(docs)
                    print(f"已加载: {file_path} ({len(docs)}个文档)")
                except Exception as e:
                    print(f"加载失败 {file_path}: {e}")
    
    return documents


def get_document_paths(*dirs: Path) -> List[str]:
    """获取所有文档路径"""
    paths = []
    
    for dir_path in dirs:
        if not dir_path.exists():
            continue
        
        for ext in SUPPORTED_EXTENSIONS:
            files = list(dir_path.rglob(f"*{ext}"))
            paths.extend([str(f) for f in files])
    
    return paths


def get_supported_extensions() -> List[str]:
    """获取支持的文件扩展名"""
    return list(SUPPORTED_EXTENSIONS)
