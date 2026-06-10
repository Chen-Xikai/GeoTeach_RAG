"""
GeoTeach AI Agent - 图片提取模块

从PPT、PDF、DOCX等文档中提取图片。
"""

import hashlib
from pathlib import Path
from typing import List, Dict, Optional


class ImageExtractor:
    """图片提取器"""
    
    def __init__(self, output_dir: str = "data/images"):
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
    
    def extract_from_pptx(self, file_path: str) -> List[dict]:
        """从PPT提取图片"""
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
        except ImportError:
            print("python-pptx未安装，请运行: pip install python-pptx")
            return []
        
        prs = Presentation(file_path)
        images = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            for shape_idx, shape in enumerate(slide.shapes):
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image = shape.image
                    image_bytes = image.blob
                    
                    # 生成唯一文件名
                    content_hash = hashlib.md5(image_bytes).hexdigest()[:8]
                    ext = self._get_image_extension(image.content_type)
                    image_name = f"pptx_{Path(file_path).stem}_{slide_num}_{shape_idx}_{content_hash}{ext}"
                    image_path = self.output_dir / image_name
                    
                    # 保存图片
                    with open(image_path, "wb") as f:
                        f.write(image_bytes)
                    
                    images.append({
                        "path": str(image_path),
                        "source": file_path,
                        "slide": slide_num,
                        "type": "image"
                    })
        
        return images
    
    def extract_from_pdf(self, file_path: str) -> List[dict]:
        """从PDF提取图片"""
        try:
            import fitz  # PyMuPDF
        except ImportError:
            print("PyMuPDF未安装，请运行: pip install PyMuPDF")
            return []
        
        doc = fitz.open(file_path)
        images = []
        
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            image_list = page.get_images(full=True)
            
            for img_idx, img in enumerate(image_list):
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                
                # 生成唯一文件名
                content_hash = hashlib.md5(image_bytes).hexdigest()[:8]
                ext = self._get_extension_from_bytes(image_bytes)
                image_name = f"pdf_{Path(file_path).stem}_{page_num+1}_{img_idx}_{content_hash}{ext}"
                image_path = self.output_dir / image_name
                
                # 保存图片
                with open(image_path, "wb") as f:
                    f.write(image_bytes)
                
                images.append({
                    "path": str(image_path),
                    "source": file_path,
                    "page": page_num + 1,
                    "type": "image"
                })
        
        doc.close()
        return images
    
    def extract_from_docx(self, file_path: str) -> List[dict]:
        """从DOCX提取图片"""
        try:
            from docx import Document
        except ImportError:
            print("python-docx未安装，请运行: pip install python-docx")
            return []
        
        doc = Document(file_path)
        images = []
        
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                image_data = rel.target_part.blob
                
                # 生成唯一文件名
                content_hash = hashlib.md5(image_data).hexdigest()[:8]
                ext = self._get_extension_from_bytes(image_data)
                image_name = f"docx_{Path(file_path).stem}_{len(images)+1}_{content_hash}{ext}"
                image_path = self.output_dir / image_name
                
                # 保存图片
                with open(image_path, "wb") as f:
                    f.write(image_data)
                
                images.append({
                    "path": str(image_path),
                    "source": file_path,
                    "type": "image"
                })
        
        return images
    
    def extract(self, file_path: str) -> List[dict]:
        """根据文件类型提取图片"""
        ext = Path(file_path).suffix.lower()
        
        if ext == '.pptx':
            return self.extract_from_pptx(file_path)
        elif ext == '.pdf':
            return self.extract_from_pdf(file_path)
        elif ext == '.docx':
            return self.extract_from_docx(file_path)
        else:
            return []
    
    def _get_image_extension(self, content_type: str) -> str:
        """根据content_type获取扩展名"""
        mapping = {
            'image/png': '.png',
            'image/jpeg': '.jpg',
            'image/gif': '.gif',
            'image/webp': '.webp',
            'image/svg+xml': '.svg'
        }
        return mapping.get(content_type, '.png')
    
    def _get_extension_from_bytes(self, image_bytes: bytes) -> str:
        """根据图片字节判断扩展名"""
        if image_bytes[:8] == b'\x89PNG\r\n\x1a\n':
            return '.png'
        elif image_bytes[:2] == b'\xff\xd8':
            return '.jpg'
        elif image_bytes[:4] == b'GIF8':
            return '.gif'
        elif image_bytes[:4] == b'RIFF' and image_bytes[8:12] == b'WEBP':
            return '.webp'
        else:
            return '.png'
