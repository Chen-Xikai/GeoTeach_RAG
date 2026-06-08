"""
GeoTeach AI Agent - 文档处理模块

整合文档读取、文本提取、多模态处理功能。
"""

import os
from pathlib import Path
from typing import List, Dict, Optional
from concurrent.futures import ThreadPoolExecutor, as_completed

from core.image_extractor import ImageExtractor
from core.ocr import VisionProcessor


# ==================== 支持的文件格式 ====================

SUPPORTED_EXTENSIONS = {'.pdf', '.docx', '.txt', '.md', '.pptx'}


# ==================== 文件读取 ====================

def read_file(filepath: str) -> str:
    """读取文件内容"""
    ext = Path(filepath).suffix.lower()
    
    if ext == '.pptx':
        return _read_pptx(filepath)
    elif ext == '.pdf':
        return _read_pdf(filepath)
    elif ext == '.docx':
        return _read_docx(filepath)
    elif ext in ['.txt', '.md']:
        return _read_text(filepath)
    else:
        raise ValueError(f"不支持的文件格式: {ext}")


def _read_pptx(filepath: str) -> str:
    """读取PPT文件"""
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("python-pptx未安装，请运行: pip install python-pptx")
    
    prs = Presentation(filepath)
    texts = []
    
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = [f"--- 第{slide_num}页 ---"]
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text)
            
            if shape.has_table:
                table = shape.table
                rows = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    rows.append(" | ".join(cells))
                slide_texts.append("\n".join(rows))
        
        if len(slide_texts) > 1:
            texts.append("\n".join(slide_texts))
    
    return "\n\n".join(texts)


def _read_pdf(filepath: str) -> str:
    """读取PDF文件"""
    try:
        from pypdf import PdfReader
        reader = PdfReader(filepath)
        texts = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                texts.append(text)
        return "\n\n".join(texts)
    except Exception as e:
        raise ValueError(f"PDF读取失败: {e}")


def _read_docx(filepath: str) -> str:
    """读取DOCX文件"""
    try:
        from docx import Document
        doc = Document(filepath)
        
        texts = []
        for para in doc.paragraphs:
            if para.text.strip():
                texts.append(para.text)
        
        for table in doc.tables:
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(" | ".join(cells))
            texts.append("\n".join(rows))
        
        return "\n".join(texts)
    except Exception as e:
        raise ValueError(f"DOCX读取失败: {e}")


def _read_text(filepath: str) -> str:
    """读取文本文件"""
    try:
        with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()
    except UnicodeDecodeError:
        try:
            with open(filepath, 'r', encoding='gbk', errors='ignore') as f:
                return f.read()
        except Exception:
            with open(filepath, 'r', encoding='latin-1', errors='ignore') as f:
                return f.read()


def _clean_text(text: str) -> str:
    """清理文本中的无效Unicode字符"""
    return text.encode('utf-8', errors='ignore').decode('utf-8', errors='ignore')


# ==================== 文档加载 ====================

def load_single_document(filepath: str, use_multimodal: bool = False) -> List[dict]:
    """
    加载单个文档
    
    Args:
        filepath: 文件路径
        use_multimodal: 是否使用多模态处理（提取图片、OCR等）
    
    Returns:
        文档列表
    """
    filepath = str(filepath)
    ext = Path(filepath).suffix.lower()
    
    if ext not in SUPPORTED_EXTENSIONS:
        raise ValueError(f"不支持的文件格式: {ext}")
    
    content = read_file(filepath)
    content = _clean_text(content)
    
    if not content.strip():
        return []
    
    doc = {
        "page_content": content,
        "metadata": {
            "source": filepath,
            "filename": Path(filepath).name,
            "file_type": ext,
        }
    }
    
    if use_multimodal and ext in ['.pptx', '.pdf', '.docx']:
        try:
            processor = MultimodalProcessor(use_vision_api=True)
            docs = processor.process_document(filepath)
            if docs:
                return docs
        except Exception as e:
            print(f"多模态处理失败，使用普通模式: {e}")
    
    return [doc]


def load_all_documents(*dirs: Path, use_multimodal: bool = False) -> List[dict]:
    """加载所有目录下的文档"""
    documents = []
    
    for dir_path in dirs:
        if not dir_path.exists():
            continue
        
        for ext in SUPPORTED_EXTENSIONS:
            files = list(dir_path.rglob(f"*{ext}"))
            for file_path in files:
                try:
                    docs = load_single_document(str(file_path), use_multimodal=use_multimodal)
                    documents.extend(docs)
                    print(f"已加载: {file_path} ({len(docs)}个文档)")
                except Exception as e:
                    print(f"加载失败 {file_path}: {e}")
    
    return documents


def get_document_paths(*dirs: Path) -> List[str]:
    """获取所有文档路径"""
    paths = []
    
    for dir_path in dirs:
        if not dir_path.exists():
            continue
        
        for ext in SUPPORTED_EXTENSIONS:
            files = list(dir_path.rglob(f"*{ext}"))
            paths.extend([str(f) for f in files])
    
    return paths


def get_supported_extensions() -> List[str]:
    """获取支持的文件扩展名"""
    return list(SUPPORTED_EXTENSIONS)


# ==================== 多模态处理器 ====================

class MultimodalProcessor:
    """多模态文档处理器"""
    
    def __init__(self, use_vision_api: bool = True):
        """
        初始化多模态处理器
        
        Args:
            use_vision_api: 是否使用视觉模型API
        """
        self.image_extractor = ImageExtractor()
        self.use_vision_api = use_vision_api
        
        if use_vision_api:
            try:
                self.vision_processor = VisionProcessor()
            except Exception as e:
                print(f"视觉模型初始化失败: {e}")
                self.use_vision_api = False
                self.vision_processor = None
        else:
            self.vision_processor = None
    
    def process_document(self, file_path: str) -> List[dict]:
        """
        处理包含图片的文档
        
        Args:
            file_path: 文档路径
            
        Returns:
            处理后的文档列表
        """
        documents = []
        file_path = str(file_path)
        ext = Path(file_path).suffix.lower()
        
        text_docs = self._extract_text(file_path, ext)
        documents.extend(text_docs)
        
        images = self.image_extractor.extract(file_path)
        
        if not images:
            return documents
        
        for img in images:
            img_doc = self._process_image(img)
            if img_doc:
                documents.append(img_doc)
        
        return documents
    
    def _extract_text(self, file_path: str, ext: str) -> List[dict]:
        """提取文档文本内容"""
        try:
            if ext == '.pptx':
                return self._extract_from_pptx(file_path)
            elif ext == '.pdf':
                return self._extract_from_pdf(file_path)
            elif ext == '.docx':
                return self._extract_from_docx(file_path)
            elif ext in ['.txt', '.md']:
                return self._extract_from_text(file_path)
            else:
                return []
        except Exception as e:
            print(f"文本提取失败: {e}")
            return []
    
    def _extract_from_pptx(self, file_path: str) -> List[dict]:
        """从PPT提取文本"""
        try:
            from pptx import Presentation
        except ImportError:
            return []
        
        prs = Presentation(file_path)
        documents = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            text_content = []
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_content.append(shape.text)
                
                if shape.has_table:
                    table = shape.table
                    rows = []
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        rows.append(" | ".join(cells))
                    text_content.append("\n".join(rows))
            
            if text_content:
                content = "\n".join(text_content)
                documents.append({
                    "page_content": content,
                    "metadata": {
                        "source": file_path,
                        "type": "text",
                        "slide_number": slide_num,
                        "total_slides": len(prs.slides)
                    }
                })
        
        return documents
    
    def _extract_from_pdf(self, file_path: str) -> List[dict]:
        """从PDF提取文本"""
        try:
            from langchain_community.document_loaders import PyPDFLoader
            loader = PyPDFLoader(file_path)
            docs = loader.load()
            
            documents = []
            for i, doc in enumerate(docs):
                documents.append({
                    "page_content": doc.page_content,
                    "metadata": {
                        "source": file_path,
                        "type": "text",
                        "page": i + 1
                    }
                })
            return documents
        except Exception as e:
            print(f"PDF提取失败: {e}")
            return []
    
    def _extract_from_docx(self, file_path: str) -> List[dict]:
        """从DOCX提取文本"""
        try:
            from docx import Document
            doc = Document(file_path)
            
            text_content = []
            for para in doc.paragraphs:
                if para.text.strip():
                    text_content.append(para.text)
            
            for table in doc.tables:
                rows = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    rows.append(" | ".join(cells))
                text_content.append("\n".join(rows))
            
            if text_content:
                return [{
                    "page_content": "\n".join(text_content),
                    "metadata": {
                        "source": file_path,
                        "type": "text"
                    }
                }]
            return []
        except Exception as e:
            print(f"DOCX提取失败: {e}")
            return []
    
    def _extract_from_text(self, file_path: str) -> List[dict]:
        """从文本文件提取内容"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            if content.strip():
                return [{
                    "page_content": content,
                    "metadata": {
                        "source": file_path,
                        "type": "text"
                    }
                }]
            return []
        except Exception as e:
            print(f"文本提取失败: {e}")
            return []
    
    def _process_image(self, img_info: dict) -> Optional[dict]:
        """处理单个图片"""
        image_path = img_info["path"]
        
        try:
            ocr_text = ""
            description = ""
            table_data = {}
            geo_info = {}
            
            if self.use_vision_api and self.vision_processor:
                try:
                    ocr_text = self.vision_processor.extract_text(image_path)
                    description = self.vision_processor.describe_image(image_path)
                    table_data = self.vision_processor.analyze_table(image_path)
                    geo_info = self.vision_processor.extract_geographic_info(image_path)
                except Exception as e:
                    print(f"视觉模型处理失败: {e}")
            
            content_parts = []
            
            if description:
                content_parts.append(f"【图片描述】\n{description}")
            
            if ocr_text:
                content_parts.append(f"【OCR文字】\n{ocr_text}")
            
            if table_data.get("has_table"):
                content_parts.append(f"【表格内容】\n{table_data.get('table_description', '')}")
            
            if geo_info.get("topics"):
                content_parts.append(f"【地理主题】\n{', '.join(geo_info['topics'])}")
            
            if not content_parts:
                return None
            
            content = "\n\n".join(content_parts)
            
            metadata = {
                "source": img_info.get("source", ""),
                "type": "image",
                "image_path": image_path,
                "has_ocr": bool(ocr_text),
                "has_table": table_data.get("has_table", False)
            }
            
            if "slide" in img_info:
                metadata["slide_number"] = img_info["slide"]
            if "page" in img_info:
                metadata["page"] = img_info["page"]
            
            return {
                "page_content": content,
                "metadata": metadata
            }
        
        except Exception as e:
            print(f"图片处理失败 {image_path}: {e}")
            return None
    
    def process_batch(self, file_paths: List[str], max_workers: int = 2) -> List[dict]:
        """
        批量处理文档
        
        Args:
            file_paths: 文件路径列表
            max_workers: 最大并发数
            
        Returns:
            处理后的文档列表
        """
        all_documents = []
        
        with ThreadPoolExecutor(max_workers=max_workers) as executor:
            future_to_path = {
                executor.submit(self.process_document, fp): fp 
                for fp in file_paths
            }
            for future in as_completed(future_to_path):
                file_path = future_to_path[future]
                try:
                    docs = future.result()
                    all_documents.extend(docs)
                    print(f"处理完成: {file_path} ({len(docs)}个文档)")
                except Exception as e:
                    print(f"处理失败 {file_path}: {e}")
        
        return all_documents


# ==================== 便捷函数 ====================

def process_document_multimodal(file_path: str, use_vision_api: bool = True) -> List[dict]:
    """便捷函数：处理单个文档"""
    processor = MultimodalProcessor(use_vision_api=use_vision_api)
    return processor.process_document(file_path)
